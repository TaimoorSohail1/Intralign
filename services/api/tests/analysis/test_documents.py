from io import BytesIO

import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from oslo_api.analysis.documents import (
    MAX_DOCUMENT_BYTES,
    DocumentRejected,
    OcrResult,
    parse_document,
)


def _single_page_pdf(text: str) -> bytes:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
        ),
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    output = BytesIO()
    output.write(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{index} 0 obj\n".encode() + body + b"\nendobj\n")
    xref = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode())
    output.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.write(f"{offset:010d} 00000 n \n".encode())
    output.write(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n"
        ).encode()
    )
    return output.getvalue()


def test_pdf_is_extracted_into_page_addressable_fragments() -> None:
    parsed = parse_document(
        file_name="plan.pdf",
        declared_content_type="application/pdf",
        content=_single_page_pdf("Timeline is 6 months. Budget is $1.8M."),
    )

    assert parsed.detected_content_type == "application/pdf"
    assert parsed.page_count == 1
    assert "Timeline is 6 months" in parsed.fragments[0].content
    assert parsed.fragments[0].locator["page"] == 1
    assert parsed.fragments[0].ordinal == 0


def test_docx_is_extracted_into_section_addressable_fragments() -> None:
    document = Document()
    document.add_heading("Delivery plan", level=1)
    document.add_paragraph("The approved budget is USD 2.4M.")
    content = BytesIO()
    document.save(content)

    parsed = parse_document(
        file_name="plan.docx",
        declared_content_type=(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
        content=content.getvalue(),
    )

    assert parsed.detected_content_type.endswith("wordprocessingml.document")
    assert parsed.fragments[0].content == "Delivery plan\nThe approved budget is USD 2.4M."
    assert parsed.fragments[0].locator == {
        "kind": "docx_section",
        "section": "Delivery plan",
        "paragraph_start": 1,
        "paragraph_end": 2,
    }


def test_pptx_is_extracted_into_slide_addressable_fragments() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title.text = "Program risks"
    detail = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    detail.text = "Vendor selection is unresolved."
    content = BytesIO()
    presentation.save(content)

    parsed = parse_document(
        file_name="plan.pptx",
        declared_content_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        content=content.getvalue(),
    )

    assert parsed.detected_content_type.endswith("presentationml.presentation")
    assert parsed.fragments[0].content == "Program risks\nVendor selection is unresolved."
    assert parsed.fragments[0].locator == {
        "kind": "pptx_slide",
        "slide": 1,
        "shape_start": 1,
        "shape_end": 2,
    }


def test_xlsx_is_extracted_into_sheet_and_cell_addressable_fragments() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Risks"
    sheet.append(["Risk", "Owner"])
    sheet.append(["Vendor unresolved", "Procurement"])
    content = BytesIO()
    workbook.save(content)

    parsed = parse_document(
        file_name="plan.xlsx",
        declared_content_type=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        content=content.getvalue(),
    )

    assert parsed.detected_content_type.endswith("spreadsheetml.sheet")
    assert parsed.fragments[0].content == ("Risk | Owner\nVendor unresolved | Procurement")
    assert parsed.fragments[0].locator == {
        "kind": "xlsx_range",
        "sheet": "Risks",
        "cell_range": "A1:B2",
    }


def test_scanned_pdf_uses_ocr_and_preserves_page_reference() -> None:
    image = Image.new("RGB", (600, 800), "white")
    content = BytesIO()
    image.save(content, format="PDF")

    class StubOcr:
        def extract(self, image: Image.Image) -> OcrResult:
            return OcrResult(
                text="Approved timeline is 12 months.",
                confidence=0.92,
            )

    parsed = parse_document(
        file_name="scan.pdf",
        declared_content_type="application/pdf",
        content=content.getvalue(),
        ocr_service=StubOcr(),
    )

    assert parsed.fragments[0].content == "Approved timeline is 12 months."
    assert parsed.fragments[0].locator["page"] == 1
    assert parsed.fragments[0].locator["ocr"] is True
    assert parsed.fragments[0].locator["ocr_confidence"] == 0.92


@pytest.mark.parametrize(
    ("file_name", "content_type", "content", "error_code"),
    [
        ("broken.pdf", "application/pdf", b"%PDF-not-valid", "DOCUMENT_PDF_INVALID"),
        ("plan.exe", "application/octet-stream", b"not a document", "DOCUMENT_TYPE_UNSUPPORTED"),
    ],
)
def test_invalid_documents_are_rejected_safely(
    file_name: str,
    content_type: str,
    content: bytes,
    error_code: str,
) -> None:
    with pytest.raises(DocumentRejected, match=error_code):
        parse_document(
            file_name=file_name,
            declared_content_type=content_type,
            content=content,
        )


def test_oversized_document_is_rejected_before_parsing() -> None:
    with pytest.raises(DocumentRejected, match="DOCUMENT_TOO_LARGE"):
        parse_document(
            file_name="large.txt",
            declared_content_type="text/plain",
            content=b"x" * (MAX_DOCUMENT_BYTES + 1),
        )


def test_password_protected_pdf_has_a_specific_error() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.encrypt("secret")
    content = BytesIO()
    writer.write(content)

    with pytest.raises(DocumentRejected, match="DOCUMENT_PASSWORD_PROTECTED"):
        parse_document(
            file_name="locked.pdf",
            declared_content_type="application/pdf",
            content=content.getvalue(),
        )


@pytest.mark.parametrize("file_name", ["locked.docx", "locked.pptx", "locked.xlsx"])
def test_encrypted_office_container_has_a_specific_error(file_name: str) -> None:
    with pytest.raises(DocumentRejected, match="DOCUMENT_PASSWORD_PROTECTED"):
        parse_document(
            file_name=file_name,
            declared_content_type="application/octet-stream",
            content=b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1encrypted-office",
        )


def test_file_signature_cannot_be_hidden_behind_an_unsupported_extension() -> None:
    with pytest.raises(DocumentRejected, match="DOCUMENT_TYPE_MISMATCH"):
        parse_document(
            file_name="payload.exe",
            declared_content_type="application/octet-stream",
            content=_single_page_pdf("A valid PDF with an invalid name"),
        )
