from io import BytesIO
from uuid import UUID, uuid4

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import create_engine, text

from oslo_api.analysis.document_store import DatabaseDocumentStore
from oslo_api.analysis.documents import (
    DocumentRejected,
    ParsedDocument,
    ParsedFragment,
)
from oslo_api.analysis.harness import DeterministicAgentHarness
from oslo_api.analysis.models import AnalysisRunRequest, AnalysisRunStatus, RunKind
from oslo_api.analysis.persistence import DatabaseAnalysisStore
from oslo_api.analysis.workflow import AnalysisWorkflow
from oslo_api.settings import Settings

SETTINGS = Settings()  # type: ignore[call-arg]
WORKSPACE_ID = UUID("018f9f7e-8de2-7000-8000-000000000010")


def test_uploaded_document_is_stored_and_fragmented_for_analysis(tmp_path) -> None:
    engine = create_engine(SETTINGS.database_url)
    project_id = uuid4()
    with engine.begin() as connection:
        owner_id = connection.execute(
            text("select id from auth.users where email = 'admin@oslo.local'")
        ).scalar_one()
        connection.execute(
            text(
                """
                insert into public.projects (id, workspace_id, name, status, created_by)
                values (:id, :workspace_id, 'Document ingestion', 'draft', :owner_id)
                """
            ),
            {"id": project_id, "workspace_id": WORKSPACE_ID, "owner_id": owner_id},
        )
    try:
        uploaded = DatabaseDocumentStore(engine=engine, object_root=tmp_path).ingest(
            workspace_id=WORKSPACE_ID,
            project_id=project_id,
            submitted_by=owner_id,
            file_name="plan.txt",
            declared_content_type="text/plain",
            content=b"Timeline is 6 months. Budget is $1.8M.",
        )

        assert uploaded.status == "parsed"
        assert uploaded.fragment_count == 1
        assert (tmp_path / uploaded.object_key).read_bytes().startswith(b"Timeline")
        with engine.connect() as connection:
            document = (
                connection.execute(
                    text(
                        """
                        select status, parser_version
                        from public.source_documents where id = :document_id
                        """
                    ),
                    {"document_id": uploaded.id},
                )
                .mappings()
                .one()
            )
            fragment = (
                connection.execute(
                    text(
                        """
                    select content, locator from public.source_fragments
                    where source_document_id = :document_id
                    """
                    ),
                    {"document_id": uploaded.id},
                )
                .mappings()
                .one()
            )

        assert document == {"status": "parsed", "parser_version": "oslo-parser-v2"}
        assert "Budget is $1.8M" in fragment["content"]
        assert fragment["locator"]["page"] == 1

        repeated = DatabaseDocumentStore(engine=engine, object_root=tmp_path).ingest(
            workspace_id=WORKSPACE_ID,
            project_id=project_id,
            submitted_by=owner_id,
            file_name="plan.txt",
            declared_content_type="text/plain",
            content=b"Timeline is 6 months. Budget is $1.8M.",
        )
        assert repeated.id == uploaded.id
        with engine.connect() as connection:
            assert (
                connection.execute(
                    text(
                        """
                        select count(*) from public.source_documents
                        where project_id = :project_id
                        """
                    ),
                    {"project_id": project_id},
                ).scalar_one()
                == 1
            )
    finally:
        with engine.begin() as connection:
            connection.execute(
                text("delete from public.projects where id = :id"),
                {"id": project_id},
            )


def test_analysis_reads_persisted_document_evidence(tmp_path) -> None:
    engine = create_engine(SETTINGS.database_url)
    project_id = uuid4()
    with engine.begin() as connection:
        owner_id = connection.execute(
            text("select id from auth.users where email = 'admin@oslo.local'")
        ).scalar_one()
        connection.execute(
            text(
                """
                insert into public.projects (id, workspace_id, name, status, created_by)
                values (:id, :workspace_id, 'Evidence analysis', 'draft', :owner_id)
                """
            ),
            {"id": project_id, "workspace_id": WORKSPACE_ID, "owner_id": owner_id},
        )
    try:
        uploaded = DatabaseDocumentStore(engine=engine, object_root=tmp_path).ingest(
            workspace_id=WORKSPACE_ID,
            project_id=project_id,
            submitted_by=owner_id,
            file_name="stress.txt",
            declared_content_type="text/plain",
            content=(
                b"Initial duration is 9 months, another section states 6 months, "
                b"and the roadmap states 12 months. Budget is $2.5M, $2.1M or $1.8M. "
                b"Mobile app, HR module and inventory integration are inconsistently included. "
                b"Success metrics are not defined. Production deployment requires regulatory "
                b"approval not included in the timeline. Resource allocations conflict. "
                b"Final vendor undecided. Migration volume unknown."
            ),
        )
        store = DatabaseAnalysisStore(engine)
        result = AnalysisWorkflow(
            store=store,
            harness=DeterministicAgentHarness(),
        ).run(
            AnalysisRunRequest(
                workspace_id=WORKSPACE_ID,
                project_id=project_id,
                requested_by=owner_id,
                kind=RunKind.EXTENDED,
                description="",
                source_names=("stress.txt",),
                source_document_ids=(uploaded.id,),
                idempotency_key=f"document-evidence:{project_id}",
            )
        )

        assert result.status is AnalysisRunStatus.COMPLETED
        assert result.snapshot is not None
        assert result.snapshot.assessment.confidence_band == "Low"
        titles = {issue.title for issue in result.snapshot.assessment.issues}
        assert "Conflicting project timelines" in titles
        assert "Deployment depends on unresolved regulatory approval" in titles
        assert all(
            reference.startswith(f"document:{uploaded.id}:")
            for issue in result.snapshot.assessment.issues
            for reference in issue.evidence_refs
        )
        with engine.connect() as connection:
            assert (
                connection.execute(
                    text(
                        """
                    select count(*) from public.analysis_node_attempts
                    where analysis_run_id = :run_id and status = 'completed'
                    """
                    ),
                    {"run_id": result.run_id},
                ).scalar_one()
                == 12
            )
    finally:
        with engine.begin() as connection:
            connection.execute(
                text("delete from public.projects where id = :id"),
                {"id": project_id},
            )


def test_mixed_office_formats_share_traceable_evidence_and_one_workflow(tmp_path) -> None:
    engine = create_engine(SETTINGS.database_url)
    project_id = uuid4()
    with engine.begin() as connection:
        owner_id = connection.execute(
            text("select id from auth.users where email = 'admin@oslo.local'")
        ).scalar_one()
        connection.execute(
            text(
                """
                insert into public.projects (id, workspace_id, name, status, created_by)
                values (:id, :workspace_id, 'Mixed evidence', 'draft', :owner_id)
                """
            ),
            {"id": project_id, "workspace_id": WORKSPACE_ID, "owner_id": owner_id},
        )

    document = Document()
    document.add_heading("Scope", level=1)
    document.add_paragraph("Phase one includes customer migration.")
    docx = BytesIO()
    document.save(docx)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    textbox.text = "Vendor selection is unresolved."
    pptx = BytesIO()
    presentation.save(pptx)

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Scenario", "Amount"])
    sheet.append(["Approved", "Pending"])
    xlsx = BytesIO()
    workbook.save(xlsx)

    sources = (
        ("scope.docx", docx.getvalue()),
        ("dependencies.pptx", pptx.getvalue()),
        ("budget.xlsx", xlsx.getvalue()),
    )
    try:
        document_store = DatabaseDocumentStore(engine=engine, object_root=tmp_path)
        uploaded = tuple(
            document_store.ingest(
                workspace_id=WORKSPACE_ID,
                project_id=project_id,
                submitted_by=owner_id,
                file_name=file_name,
                declared_content_type="application/octet-stream",
                content=content,
            )
            for file_name, content in sources
        )
        store = DatabaseAnalysisStore(engine)
        request = AnalysisRunRequest(
            workspace_id=WORKSPACE_ID,
            project_id=project_id,
            requested_by=owner_id,
            kind=RunKind.INITIAL,
            description="",
            source_names=tuple(file_name for file_name, _ in sources),
            source_document_ids=tuple(item.id for item in uploaded),
            idempotency_key=f"mixed-formats:{project_id}",
        )

        evidence_refs = tuple(item.reference for item in store.evidence_for(request))
        result = AnalysisWorkflow(
            store=store,
            harness=DeterministicAgentHarness(),
        ).run(request)

        assert any(":section:Scope:" in reference for reference in evidence_refs)
        assert any(":slide:1:" in reference for reference in evidence_refs)
        assert any(":sheet:Budget:range:A1%3AB2:" in reference for reference in evidence_refs)
        assert result.status is AnalysisRunStatus.COMPLETED
        with engine.connect() as connection:
            assert (
                connection.execute(
                    text(
                        """
                    select count(*) from public.analysis_node_attempts
                    where analysis_run_id = :run_id and status = 'completed'
                    """
                    ),
                    {"run_id": result.run_id},
                ).scalar_one()
                == 12
            )
    finally:
        with engine.begin() as connection:
            connection.execute(
                text("delete from public.projects where id = :id"),
                {"id": project_id},
            )


def test_transient_parser_failure_retries_and_records_attempts(tmp_path) -> None:
    engine = create_engine(SETTINGS.database_url)
    project_id = uuid4()
    with engine.begin() as connection:
        owner_id = connection.execute(
            text("select id from auth.users where email = 'admin@oslo.local'")
        ).scalar_one()
        connection.execute(
            text(
                """
                insert into public.projects (id, workspace_id, name, status, created_by)
                values (:id, :workspace_id, 'Retry parsing', 'draft', :owner_id)
                """
            ),
            {"id": project_id, "workspace_id": WORKSPACE_ID, "owner_id": owner_id},
        )

    class TransientParser:
        def __init__(self) -> None:
            self.calls = 0

        def __call__(self, **_) -> ParsedDocument:
            self.calls += 1
            if self.calls < 3:
                raise DocumentRejected("DOCUMENT_OCR_FAILED")
            return ParsedDocument(
                detected_content_type="application/pdf",
                page_count=1,
                fragments=(
                    ParsedFragment(
                        ordinal=0,
                        content="Recovered evidence",
                        locator={"kind": "pdf_page", "page": 1},
                    ),
                ),
            )

    parser = TransientParser()
    try:
        uploaded = DatabaseDocumentStore(
            engine=engine,
            object_root=tmp_path,
            parser=parser,
            max_parse_attempts=3,
        ).ingest(
            workspace_id=WORKSPACE_ID,
            project_id=project_id,
            submitted_by=owner_id,
            file_name="scan.pdf",
            declared_content_type="application/pdf",
            content=b"%PDF-retryable",
        )

        assert uploaded.status == "parsed"
        assert parser.calls == 3
        with engine.connect() as connection:
            document = (
                connection.execute(
                    text(
                        """
                        select status, parse_attempt_count
                        from public.source_documents where id = :document_id
                        """
                    ),
                    {"document_id": uploaded.id},
                )
                .mappings()
                .one()
            )
            attempts = connection.execute(
                text(
                    """
                    select status, error_code
                    from public.document_parse_attempts
                    where source_document_id = :document_id
                    order by attempt_no
                    """
                ),
                {"document_id": uploaded.id},
            ).all()

        assert document == {"status": "parsed", "parse_attempt_count": 3}
        assert attempts == [
            ("failed", "DOCUMENT_OCR_FAILED"),
            ("failed", "DOCUMENT_OCR_FAILED"),
            ("completed", None),
        ]
    finally:
        with engine.begin() as connection:
            connection.execute(
                text("delete from public.projects where id = :id"),
                {"id": project_id},
            )


def test_reupload_resumes_a_failed_document_without_creating_a_duplicate(tmp_path) -> None:
    engine = create_engine(SETTINGS.database_url)
    project_id = uuid4()
    with engine.begin() as connection:
        owner_id = connection.execute(
            text("select id from auth.users where email = 'admin@oslo.local'")
        ).scalar_one()
        connection.execute(
            text(
                """
                insert into public.projects (id, workspace_id, name, status, created_by)
                values (:id, :workspace_id, 'Resume parsing', 'draft', :owner_id)
                """
            ),
            {"id": project_id, "workspace_id": WORKSPACE_ID, "owner_id": owner_id},
        )

    def failing_parser(**_) -> ParsedDocument:
        raise DocumentRejected("DOCUMENT_OCR_FAILED")

    def recovered_parser(**_) -> ParsedDocument:
        return ParsedDocument(
            detected_content_type="application/pdf",
            page_count=1,
            fragments=(
                ParsedFragment(
                    ordinal=0,
                    content="Recovered after re-upload",
                    locator={"kind": "pdf_page", "page": 1},
                ),
            ),
        )

    try:
        with pytest.raises(DocumentRejected, match="DOCUMENT_OCR_FAILED"):
            DatabaseDocumentStore(
                engine=engine,
                object_root=tmp_path,
                parser=failing_parser,
                max_parse_attempts=2,
            ).ingest(
                workspace_id=WORKSPACE_ID,
                project_id=project_id,
                submitted_by=owner_id,
                file_name="scan.pdf",
                declared_content_type="application/pdf",
                content=b"%PDF-resumable",
            )

        uploaded = DatabaseDocumentStore(
            engine=engine,
            object_root=tmp_path,
            parser=recovered_parser,
            max_parse_attempts=2,
        ).ingest(
            workspace_id=WORKSPACE_ID,
            project_id=project_id,
            submitted_by=owner_id,
            file_name="scan.pdf",
            declared_content_type="application/pdf",
            content=b"%PDF-resumable",
        )

        assert uploaded.status == "parsed"
        assert uploaded.fragment_count == 1
        with engine.connect() as connection:
            assert (
                connection.execute(
                    text(
                        """
                    select count(*) from public.source_documents
                    where project_id = :project_id
                    """
                    ),
                    {"project_id": project_id},
                ).scalar_one()
                == 1
            )
            attempts = connection.execute(
                text(
                    """
                    select attempt_no, status from public.document_parse_attempts
                    where source_document_id = :document_id order by attempt_no
                    """
                ),
                {"document_id": uploaded.id},
            ).all()
        assert attempts == [(1, "failed"), (2, "failed"), (3, "completed")]
    finally:
        with engine.begin() as connection:
            connection.execute(
                text("delete from public.projects where id = :id"),
                {"id": project_id},
            )
