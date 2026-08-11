from dataclasses import dataclass
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from typing import Protocol

import numpy as np
import pypdfium2 as pdfium
from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
from rapidocr_onnxruntime import RapidOCR

MAX_DOCUMENT_BYTES = 10 * 1024 * 1024
CHUNK_SIZE = 2_000
CHUNK_OVERLAP = 200
MAX_SPREADSHEET_CELLS = 200_000


class DocumentRejected(ValueError):
    """Raised when an uploaded source cannot be safely processed."""


LocatorValue = str | int | float | bool


@dataclass(frozen=True, slots=True)
class ParsedFragment:
    ordinal: int
    content: str
    locator: dict[str, LocatorValue]


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    detected_content_type: str
    page_count: int
    fragments: tuple[ParsedFragment, ...]


class DocumentParser(Protocol):
    """Converts a supported source into format-neutral, traceable fragments."""

    def supports(
        self,
        *,
        suffix: str,
        declared_content_type: str | None,
        content: bytes,
    ) -> bool: ...

    def parse(self, content: bytes) -> ParsedDocument: ...


@dataclass(frozen=True, slots=True)
class OcrResult:
    text: str
    confidence: float


class OcrService(Protocol):
    """Extracts text from one rendered source page."""

    def extract(self, image: Image.Image) -> OcrResult: ...


class RapidOcrService:
    """Local, replaceable OCR implementation backed by RapidOCR."""

    def __init__(self) -> None:
        self._engine = RapidOCR()

    def extract(self, image: Image.Image) -> OcrResult:
        result, _ = self._engine(np.asarray(image.convert("RGB")))
        if not result:
            return OcrResult(text="", confidence=0)
        lines = [str(item[1]).strip() for item in result if str(item[1]).strip()]
        scores = [float(item[2]) for item in result if len(item) > 2]
        return OcrResult(
            text="\n".join(lines),
            confidence=round(sum(scores) / len(scores), 4) if scores else 0,
        )


@lru_cache(maxsize=1)
def default_ocr_service() -> OcrService:
    return RapidOcrService()


def parse_document(
    *,
    file_name: str,
    declared_content_type: str | None,
    content: bytes,
    ocr_service: OcrService | None = None,
) -> ParsedDocument:
    if not content:
        raise DocumentRejected("DOCUMENT_EMPTY")
    if len(content) > MAX_DOCUMENT_BYTES:
        raise DocumentRejected("DOCUMENT_TOO_LARGE")

    suffix = Path(file_name).suffix.lower()
    if content.startswith(b"%PDF-"):
        if suffix != ".pdf":
            raise DocumentRejected("DOCUMENT_TYPE_MISMATCH")
        return _parse_pdf(content, ocr_service=ocr_service)
    if suffix in {".docx", ".pptx", ".xlsx"} and content.startswith(b"\xd0\xcf\x11\xe0"):
        raise DocumentRejected("DOCUMENT_PASSWORD_PROTECTED")
    if suffix == ".docx" and content.startswith(b"PK"):
        return _parse_docx(content)
    if suffix == ".pptx" and content.startswith(b"PK"):
        return _parse_pptx(content)
    if suffix == ".xlsx" and content.startswith(b"PK"):
        return _parse_xlsx(content)
    if suffix in {".txt", ".md", ".csv"} and (
        declared_content_type is None
        or declared_content_type.startswith("text/")
        or declared_content_type in {"application/octet-stream", "text/csv"}
    ):
        return _parse_text(content, suffix)
    raise DocumentRejected("DOCUMENT_TYPE_UNSUPPORTED")


def _parse_pdf(
    content: bytes,
    *,
    ocr_service: OcrService | None,
) -> ParsedDocument:
    try:
        reader = PdfReader(BytesIO(content))
        if reader.is_encrypted and reader.decrypt("") == 0:
            raise DocumentRejected("DOCUMENT_PASSWORD_PROTECTED")
        pages = tuple((page.extract_text() or "").strip() for page in reader.pages)
    except DocumentRejected:
        raise
    except Exception as error:
        raise DocumentRejected("DOCUMENT_PDF_INVALID") from error

    metadata: list[dict[str, LocatorValue]] = [{"kind": "pdf_page"} for _ in pages]
    if any(not page for page in pages):
        service = ocr_service or default_ocr_service()
        rendered_pdf = None
        try:
            rendered_pdf = pdfium.PdfDocument(content)
            ocr_pages = list(pages)
            for page_index, page_text in enumerate(ocr_pages):
                if page_text:
                    continue
                page = rendered_pdf[page_index]
                try:
                    image = page.render(scale=2).to_pil()
                    result = service.extract(image)
                finally:
                    page.close()
                ocr_pages[page_index] = result.text.strip()
                metadata[page_index] = {
                    "kind": "pdf_page",
                    "ocr": True,
                    "ocr_confidence": result.confidence,
                }
            pages = tuple(ocr_pages)
        except Exception as error:
            raise DocumentRejected("DOCUMENT_OCR_FAILED") from error
        finally:
            if rendered_pdf is not None:
                rendered_pdf.close()
    if not any(pages):
        raise DocumentRejected("DOCUMENT_TEXT_NOT_EXTRACTABLE")
    return ParsedDocument(
        detected_content_type="application/pdf",
        page_count=len(pages),
        fragments=_fragment_pages(pages, page_metadata=tuple(metadata)),
    )


def _parse_docx(content: bytes) -> ParsedDocument:
    try:
        document = Document(BytesIO(content))
    except PackageNotFoundError as error:
        raise DocumentRejected("DOCUMENT_DOCX_INVALID") from error
    except Exception as error:
        raise DocumentRejected("DOCUMENT_DOCX_INVALID") from error

    sections: list[tuple[str, int, int, list[str]]] = []
    current_title = "Document"
    start = 1
    lines: list[str] = []
    for paragraph_number, paragraph in enumerate(document.paragraphs, start=1):
        text = " ".join(paragraph.text.split())
        if not text:
            continue
        if paragraph.style.name.startswith("Heading"):
            if lines:
                sections.append((current_title, start, paragraph_number - 1, lines))
            current_title = text
            start = paragraph_number
            lines = [text]
        else:
            lines.append(text)
    if lines:
        sections.append((current_title, start, len(document.paragraphs), lines))
    if not sections:
        raise DocumentRejected("DOCUMENT_TEXT_NOT_EXTRACTABLE")

    fragments = tuple(
        ParsedFragment(
            ordinal=ordinal,
            content="\n".join(lines),
            locator={
                "kind": "docx_section",
                "section": title,
                "paragraph_start": paragraph_start,
                "paragraph_end": paragraph_end,
            },
        )
        for ordinal, (title, paragraph_start, paragraph_end, lines) in enumerate(sections)
    )
    return ParsedDocument(
        detected_content_type=(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
        page_count=len(sections),
        fragments=fragments,
    )


def _parse_pptx(content: bytes) -> ParsedDocument:
    try:
        presentation = Presentation(BytesIO(content))
    except Exception as error:
        raise DocumentRejected("DOCUMENT_PPTX_INVALID") from error

    fragments: list[ParsedFragment] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        first_shape: int | None = None
        last_shape: int | None = None
        for shape_number, shape in enumerate(slide.shapes, start=1):
            shape_lines: list[str] = []
            if getattr(shape, "has_text_frame", False):
                shape_lines.extend(
                    " ".join(paragraph.text.split())
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [" ".join(cell.text.split()) for cell in row.cells]
                    if any(values):
                        shape_lines.append(" | ".join(values))
            if shape_lines:
                first_shape = first_shape or shape_number
                last_shape = shape_number
                lines.extend(shape_lines)
        if not lines:
            continue
        fragments.append(
            ParsedFragment(
                ordinal=len(fragments),
                content="\n".join(lines),
                locator={
                    "kind": "pptx_slide",
                    "slide": slide_number,
                    "shape_start": first_shape or 1,
                    "shape_end": last_shape or 1,
                },
            )
        )
    if not fragments:
        raise DocumentRejected("DOCUMENT_TEXT_NOT_EXTRACTABLE")
    return ParsedDocument(
        detected_content_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        page_count=len(presentation.slides),
        fragments=tuple(fragments),
    )


def _parse_xlsx(content: bytes) -> ParsedDocument:
    try:
        workbook = load_workbook(
            BytesIO(content),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
    except Exception as error:
        raise DocumentRejected("DOCUMENT_XLSX_INVALID") from error

    fragments: list[ParsedFragment] = []
    cell_count = 0
    try:
        for sheet in workbook.worksheets:
            rows: list[tuple[int, list[str]]] = []
            max_column = 0
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = ["" if value is None else str(value).strip() for value in row]
                while values and not values[-1]:
                    values.pop()
                if not values or not any(values):
                    continue
                cell_count += len(values)
                if cell_count > MAX_SPREADSHEET_CELLS:
                    raise DocumentRejected("DOCUMENT_XLSX_CELL_LIMIT_EXCEEDED")
                max_column = max(max_column, len(values))
                rows.append((row_number, values))
            if not rows:
                continue
            start_row = rows[0][0]
            end_row = rows[-1][0]
            fragments.append(
                ParsedFragment(
                    ordinal=len(fragments),
                    content="\n".join(" | ".join(values) for _, values in rows),
                    locator={
                        "kind": "xlsx_range",
                        "sheet": sheet.title,
                        "cell_range": (f"A{start_row}:{get_column_letter(max_column)}{end_row}"),
                    },
                )
            )
    finally:
        workbook.close()
    if not fragments:
        raise DocumentRejected("DOCUMENT_TEXT_NOT_EXTRACTABLE")
    return ParsedDocument(
        detected_content_type=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
        page_count=len(workbook.sheetnames),
        fragments=tuple(fragments),
    )


def _parse_text(content: bytes, suffix: str) -> ParsedDocument:
    try:
        text = content.decode("utf-8-sig").strip()
    except UnicodeDecodeError as error:
        raise DocumentRejected("DOCUMENT_TEXT_ENCODING_UNSUPPORTED") from error
    if not text:
        raise DocumentRejected("DOCUMENT_TEXT_NOT_EXTRACTABLE")
    content_type = "text/csv" if suffix == ".csv" else "text/plain"
    return ParsedDocument(
        detected_content_type=content_type,
        page_count=1,
        fragments=_fragment_pages((text,)),
    )


def _fragment_pages(
    pages: tuple[str, ...],
    *,
    page_metadata: tuple[dict[str, LocatorValue], ...] | None = None,
) -> tuple[ParsedFragment, ...]:
    fragments: list[ParsedFragment] = []
    ordinal = 0
    for page_number, page in enumerate(pages, start=1):
        normalized = " ".join(page.split())
        start = 0
        while start < len(normalized):
            target_end = min(start + CHUNK_SIZE, len(normalized))
            end = target_end
            if target_end < len(normalized):
                boundary = normalized.rfind(" ", start, target_end)
                if boundary > start:
                    end = boundary
            fragment = normalized[start:end].strip()
            if fragment:
                fragments.append(
                    ParsedFragment(
                        ordinal=ordinal,
                        content=fragment,
                        locator={
                            **(
                                page_metadata[page_number - 1]
                                if page_metadata is not None
                                else {"kind": "text_page"}
                            ),
                            "page": page_number,
                            "char_start": start,
                            "char_end": end,
                        },
                    )
                )
                ordinal += 1
            if end >= len(normalized):
                break
            start = max(end - CHUNK_OVERLAP, start + 1)
    return tuple(fragments)
